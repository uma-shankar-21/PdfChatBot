# agents/ingestion_agent.py

import os
import pandas as pd
from pptx import Presentation
from docx import Document
from PyPDF2 import PdfReader

def parse_pdf(file):
    reader = PdfReader(file)
    text = ""
    for page in reader.pages:
        text += page.extract_text() + "\n"
    return text

def parse_docx(file):
    doc = Document(file)
    return "\n".join([para.text for para in doc.paragraphs])

def parse_pptx(file):
    prs = Presentation(file)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def parse_csv(file):
    df = pd.read_csv(file)
    return df.to_csv(index=False)

def parse_txt(file):
    return file.read().decode('utf-8')

def ingest(file, filename):
    ext = filename.split('.')[-1].lower()
    if ext == 'pdf':
        return parse_pdf(file)
    elif ext == 'docx':
        return parse_docx(file)
    elif ext == 'pptx':
        return parse_pptx(file)
    elif ext == 'csv':
        return parse_csv(file)
    elif ext in ['txt', 'md']:
        return parse_txt(file)
    else:
        return ""